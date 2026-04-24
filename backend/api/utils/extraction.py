"""
File content extraction utilities for BodhAI.

Supported formats:
  - PDF  (via PyMuPDF / fitz)
  - PPTX (via python-pptx)
  - PPT  (fallback: treated as raw text extraction)
  - Images: JPG, JPEG, PNG (via Groq vision LLM)
"""

import base64
import io

import fitz  # PyMuPDF
from pptx import Presentation
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage

# ── Constants ─────────────────────────────────────────────────
# Truncate extracted text to avoid overwhelming the LLM context window
MAX_CHARS = 12_000


def _truncate(text: str) -> str:
    if len(text) > MAX_CHARS:
        return text[:MAX_CHARS] + "\n\n[... content truncated for processing ...]"
    return text


# ── PDF ───────────────────────────────────────────────────────
def extract_pdf_text(file_obj) -> str:
    """Extract plain text from a PDF file object."""
    try:
        data = file_obj.read()
        doc = fitz.open(stream=io.BytesIO(data), filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return _truncate("\n".join(pages).strip())
    except Exception as exc:
        print(f"[BodhAI] PDF extraction error: {exc}")
        return ""


# ── PPTX ──────────────────────────────────────────────────────
def extract_pptx_text(file_obj) -> str:
    """Extract text from all slides in a PPTX file object."""
    try:
        prs = Presentation(file_obj)
        slide_texts = []
        for i, slide in enumerate(prs.slides, start=1):
            shapes_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_text.append(shape.text.strip())
            if shapes_text:
                slide_texts.append(f"[Slide {i}]\n" + "\n".join(shapes_text))
        return _truncate("\n\n".join(slide_texts))
    except Exception as exc:
        print(f"[BodhAI] PPTX extraction error: {exc}")
        return ""


# ── Images ────────────────────────────────────────────────────
def extract_image_description(file_obj, mime_type: str = "image/jpeg") -> str:
    """
    Use Groq's vision model to describe / OCR an image.
    Returns a text description suitable for the learning pipeline.
    """
    try:
        data = file_obj.read()
        b64 = base64.b64encode(data).decode("utf-8")
        llm = ChatGroq(model="llama-3.2-11b-vision-preview", temperature=0)
        msg = HumanMessage(content=[
            {
                "type": "text",
                "text": (
                    "You are an AI assistant helping extract learning content from an image. "
                    "Please do the following:\n"
                    "1. Read and transcribe any text visible in the image.\n"
                    "2. Describe any diagrams, charts, or visual elements in detail.\n"
                    "3. Summarize what the image is about for educational purposes.\n"
                    "Provide a thorough, structured response."
                ),
            },
            {
                "type": "image_url",
                "image_url": {"url": f"data:{mime_type};base64,{b64}"},
            },
        ])
        result = llm.invoke([msg])
        return _truncate(result.content)
    except Exception as exc:
        print(f"[BodhAI] Image vision error: {exc}")
        return "Image could not be processed."


# ── Dispatcher ────────────────────────────────────────────────
MIME_MAP = {
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png":  "image/png",
}

def extract_file_content(file_obj, filename: str) -> tuple[str, str]:
    """
    Dispatch extraction based on file extension.

    Returns:
        (extracted_text, file_type_label)
        file_type_label: 'pdf' | 'pptx' | 'image' | 'unknown'
    """
    name = filename.lower()

    if name.endswith(".pdf"):
        return extract_pdf_text(file_obj), "pdf"

    if name.endswith(".pptx") or name.endswith(".ppt"):
        return extract_pptx_text(file_obj), "pptx"

    for ext, mime in MIME_MAP.items():
        if name.endswith(ext):
            return extract_image_description(file_obj, mime), "image"

    return "", "unknown"
